import os
import glob
import json
from pptx import Presentation
import win32com.client

def extract_text_from_ppt_file(ppt_file_path):
    assert os.path.exists(ppt_file_path), "The specified ppt file path is not existing!"
    _, name = os.path.split(ppt_file_path)
    name = name.rsplit('.')[0]
    #assert os.path.exists(txt_file_folder), "the specified txt file folder is not existing!"
    ppt = Presentation(ppt_file_path)
    txt_box = [name]
    for slide in list(ppt.slides)[1:]:
        txt_temp = []
        for shape in slide.shapes:
            if shape.has_text_frame and not shape.text.startswith('启应经文'):
                txt_temp.append(shape.text.replace('\u3000','　'))
        if len(txt_temp)==3:
            txt_box.append(''.join(txt_temp[0:2]))
            txt_box.append(txt_temp[2])
        elif len(txt_temp)==1:
            txt_box.append(txt_temp[0])
        else:
            txt_box.append('\n'.join(txt_temp))

    return txt_box
    #with open(os.path.join(txt_file_folder, txt_name),'w', encoding='utf-8') as f:
    #    f.write('\n'.join(txt_box))

def extract_song_from_ppt_file(ppt_file_path = "C:\\Users\\qiucanro\\Downloads\\2023-0723_checked.pptx"):
    #return a list of list item, each list item is of form [title, albulm, song_scripts]
    assert os.path.exists(ppt_file_path), "The specified ppt file path is not existing!"
    #_, name = os.path.split(ppt_file_path)
    #name = name.rsplit('.')[0]
    #assert os.path.exists(txt_file_folder), "the specified txt file folder is not existing!"
    ppt = Presentation(ppt_file_path)
    full_list = []
    slides = list(ppt.slides)
    song_list, title_pos = get_song_list(slides)
    slides = slides[title_pos+1:]
    for song_title in song_list:
        txt_box = [song_title]
        for i, slide in enumerate(slides):
            if len(slide.shapes)>=2 and hasattr(slide.shapes[0],'text') and hasattr(slide.shapes[1],'text') and \
                    (slide.shapes[0].text.strip().startswith(song_title) or slide.shapes[1].text.strip().startswith(song_title)):
                if slide.shapes[0].text.strip().startswith(song_title):
                    txt_box.append(slide.shapes[1].text.strip())
                elif slide.shapes[1].text.strip().startswith(song_title):
                    txt_box.append(slide.shapes[1].text.strip().rsplit('\n')[1])
                move_on, j = True, 1
                while move_on:
                    if len(slides[i+j].shapes)>1:
                        if len(list(slides[i+j].shapes)[0].text.strip())>len(list(slides[i+j].shapes)[1].text.strip()):
                            txt_box.append(list(slides[i+j].shapes)[0].text.strip())
                        else:
                            txt_box.append(list(slides[i+j].shapes)[1].text.strip())
                        j+=1
                    else:
                        move_on = False
                break
        full_list.append(txt_box)
    return full_list

def get_song_list(slides):
    songs = []
    title_pos = None
    for i, slide in enumerate(slides):
        if len(slide.shapes)>=2:
            if list(slide.shapes)[0].text.strip().startswith('诗歌赞美'):
                songs = songs + list(slide.shapes)[1].text.strip().rsplit('\n')
                title_pos = i
            elif list(slide.shapes)[0].text.strip().startswith('回应诗歌'):
                songs = songs + [list(slide.shapes)[1].text.strip().rsplit('\n')[0]]
    return songs, title_pos

def save_ppt_content_to_json(folder="C:\\Users\\qiucanro\\Downloads\\Responsive Reading", json_file='scriptures.json'):
    assert os.path.exists(folder), "The specified ppt file folder is not existing!"
    results = {}
    for each in glob.glob(os.path.join(folder, '*.pptx')):
        temp = extract_text_from_ppt_file(each)
        results[temp[0].lstrip(' ').replace('\n','')] = '\n'.join(temp[1:])
    with open(os.path.join(folder, json_file),'w', encoding='utf-8') as f:
        json.dump(results, f)

def save_ppt_to_pptx_in_folder(folder):
    assert os.path.exists(folder), "The specified ppt file folder is not existing!"
    PptApp = win32com.client.Dispatch("Powerpoint.Application")
    PptApp.Visible = True
    for each in glob.glob(os.path.join(folder, '*.ppt')):
        print(each)
        PPtPresentation = PptApp.Presentations.Open(each)
        PPtPresentation.SaveAs(each+'x', 24)
        PPtPresentation.close()
        #save_ppt_as_pptx(each)
    PptApp.Quit()

def test_extract_text_from_ppt_file():
    ppt_path = "C:\\Users\\qiucanro\\Downloads\\Responsive Reading\\02慈爱广大.pptx"
    txt_folder = "C:\\Users\\qiucanro\\Downloads"
    txt_name = '02慈爱广大.txt'
    extract_text_from_ppt_file(ppt_path, txt_folder, txt_name)
